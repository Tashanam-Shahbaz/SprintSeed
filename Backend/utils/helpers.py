import os
import tempfile
from utils import logger
from typing import List, Optional , Dict, Any
from fastapi import UploadFile
import fitz  # PyMuPDF
import pandas as pd

from pptx import Presentation
from unstructured.partition.docx import partition_docx
from langchain_community.document_loaders import PyPDFium2Loader
from langchain_community.document_loaders import (
    UnstructuredExcelLoader,
    UnstructuredWordDocumentLoader,
    UnstructuredPowerPointLoader
)
import shutil


def read_text_file(file_path):
    """Read content from a text file."""
    with open(file_path, "r", encoding="utf-8") as file:
        return file.read()

def combine_docs(docs):
    """Combine multiple document pieces into a single string."""
    if isinstance(docs, list):
        return "\n".join(str(d) for d in docs)
    return str(docs)

def detect_scanned(file_path):
    """Detect if a PDF is scanned (contains mostly images instead of text)."""
    try:
        text_length = 0
        image_count = 0
        
        with fitz.open(file_path) as doc:
            for page in doc:
                text_length += len(page.get_text())
                image_count += len(page.get_images())
        
        # If there's very little text but many images, it's likely scanned
        return text_length < 100 and image_count > 0
    except Exception as e:
        logger.error(f"Error detecting if PDF is scanned: {e}")
        return False

def extract_text_from_scanned_pdf(file_path):
    """Extract text from a scanned PDF using OCR-like capabilities of PyMuPDF."""
    try:
        text = ""
        with fitz.open(file_path) as doc:
            for page in doc:
                text += page.get_text("text")  # Try basic text extraction first
                
                if not text.strip():  # If no text found, try getting text from images
                    text += page.get_text("dict")  # More detailed extraction
        
        return text
    except Exception as e:
        logger.error(f"Error extracting text from scanned PDF: {e}")
        return ""

def fallback_loader(file_path, file_type, extract_images=False):
    """Fallback method to load document content when primary loaders fail."""
    if file_type == "pdf":
        with fitz.open(file_path) as doc:
            return "\n".join(page.get_text() for page in doc)
    elif file_type in ["xlsx", "xls"]:
        excel_data = pd.ExcelFile(file_path)
        return "\n".join(excel_data.parse(sheet).to_string() for sheet in excel_data.sheet_names)
    elif file_type in ["docx", "doc"]:
        elements = partition_docx(
            filename=file_path,
            extract_images=extract_images,
            include_page_breaks=True,
            strategy='hi_res'
        )
        return "\n".join(str(e) for e in elements)
    elif file_type in ["pptx", "ppt"]:
        prs = Presentation(file_path)
        return "\n".join(
            shape.text for slide in prs.slides if hasattr(slide, "shapes")
            for shape in slide.shapes if hasattr(shape, "text")
        )
    elif file_type == "txt":
        return read_text_file(file_path)
    else:
        raise ValueError(f"Unsupported file type for fallback: {file_type}")

def get_file_text(file_path="", extract_images=False, file_type=None):
    """Extract text from various document formats."""
    try:
        if not file_type:
            file_type = file_path.split(".")[-1].lower()
        
        logger.info(f"Processing file: {file_path} of type: {file_type}")
        
        if file_type == "pdf":
            if detect_scanned(file_path):
                text = extract_text_from_scanned_pdf(file_path)
                logger.info(f"Detected scanned PDF: {file_path}. Content length: {len(text)}")
                if text:
                    return text
               
            docs = PyPDFium2Loader(file_path, extract_images=extract_images).load()
        elif file_type in ["xlsx", "xls"]:
            docs = UnstructuredExcelLoader(file_path, mode="elements").load()
        elif file_type in ["docx", "doc"]:
            docs = UnstructuredWordDocumentLoader(file_path).load()
        elif file_type in ["pptx", "ppt"]:
            docs = UnstructuredPowerPointLoader(file_path).load()
        elif file_type == "txt":
            docs = [read_text_file(file_path)]
        else:
            raise ValueError(f"Unsupported file type: {file_type}")
        
        return combine_docs(docs)
    except Exception as primary_error:
        logger.warning(f"Primary loader failed for {file_path}: {primary_error}. Attempting fallback loaders...")
        try:
            return fallback_loader(file_path, file_type, extract_images)
        except Exception as fallback_error:
            logger.error(f"Fallback loader failed for {file_path}: {fallback_error}")
            return ''

def process_files_for_storage(files: List[UploadFile]) -> List[Dict[str, Any]]:
    """
    Process uploaded files and prepare them for database storage.
    Returns a list of dictionaries with file information.
    """
    if not files:
        return []
    
    processed_files = []
    
    for file in files:
        try:
            # Create a temporary file to store the uploaded content
            with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file.filename.split('.')[-1]}") as temp_file:
                temp_file_path = temp_file.name
                
            # Write the uploaded file content to the temporary file using shutil
            with open(temp_file_path, "wb") as buffer:
                shutil.copyfileobj(file.file, buffer)
            
            # Process the file based on its type
            file_type = file.filename.split('.')[-1].lower()
            
            # Get the file size
            file_size = os.path.getsize(temp_file_path)
            
            # Read the file content as binary
            with open(temp_file_path, "rb") as f:
                file_binary = f.read()
            
            # Get the text content for display
            try:
                file_text = get_file_text(temp_file_path, extract_images=False, file_type=file_type)
            except Exception as text_error:
                logger.warning(f"Could not extract text from {file.filename}: {str(text_error)}")
                file_text = "[Content extraction not supported for this file type]"
            
            # Add file information to the list
            processed_files.append({
                "file_name": file.filename,
                "file_type": file.content_type or f"application/{file_type}",
                "file_size": file_size,
                "file_content": file_binary,  # Binary content for storage
                "file_text": file_text        # Text content for display
            })
            
            # Clean up the temporary file
            os.unlink(temp_file_path)
            
            # Reset file pointer for potential reuse
            file.file.seek(0)
            
        except Exception as e:
            logger.error(f"Error processing file {file.filename}: {str(e)}")
    
    return processed_files